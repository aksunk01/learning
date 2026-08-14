from pptx import Presentation

from .base import DocumentElement, DocumentParser, ParsedDocument

class PPTXParser(DocumentParser):

    def parse(self, file_path: str) -> ParsedDocument:
        presentation = Presentation(file_path)

        elements: list[DocumentElement] = []

        for slide_index, slide in enumerate(presentation.slides):
            text_parts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text = shape.text.strip()

                if text:
                    text_parts.append(text)

            slide_text = "\n\n".join(text_parts)

            if not slide_text:
                continue
            
            element = DocumentElement(
                text = slide_text,
                slide_number = slide_index+1,
                metadata=self._get_element_metadata(
                    file_path
                )
            )

            elements.append(element)

        return ParsedDocument(
            elements=elements,
            metadata=self._get_file_metadata(
                file_path,
                "pptx"
            )
        )